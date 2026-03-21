"""Tests for binary document parsers (Batch 2).

Strategy:
- Fallback tests (no optional dep): always pass — mock ImportError via monkeypatch or use
  direct import in an isolated way.
- Full-parse tests: skip when the library is not installed using pytest.importorskip.
- Registry presence tests: importing the module side-effects registration.
"""
from __future__ import annotations

import sys
from pathlib import Path
from unittest.mock import MagicMock, patch

import pytest

# ---------------------------------------------------------------------------
# Helpers
# ---------------------------------------------------------------------------

def _make_mock_import_error(module_name: str):
    """Return a side_effect callable that raises ImportError for `module_name`."""
    real_import = __builtins__.__import__ if hasattr(__builtins__, "__import__") else __import__

    def _import(name, *args, **kwargs):
        if name == module_name:
            raise ImportError(f"Mocked: {module_name} not installed")
        return real_import(name, *args, **kwargs)

    return _import


# ---------------------------------------------------------------------------
# PDF Parser
# ---------------------------------------------------------------------------

class TestPdfParser:
    def test_registers_in_registry(self):
        import homie_core.rag.parsers.pdf  # noqa: F401
        from homie_core.rag.parsers import PARSER_REGISTRY
        assert "pdf" in PARSER_REGISTRY

    def test_fallback_missing_dep(self, tmp_path):
        """When fitz (PyMuPDF) is not available, returns error ParsedDocument."""
        f = tmp_path / "test.pdf"
        f.write_bytes(b"%PDF-1.4 fake")

        # Remove fitz from sys.modules and block re-import
        saved = sys.modules.pop("fitz", None)
        try:
            with patch.dict(sys.modules, {"fitz": None}):
                # Force re-import of the parser module with fitz blocked
                import importlib
                import homie_core.rag.parsers.pdf as pdf_mod
                importlib.reload(pdf_mod)
                from homie_core.rag.parsers import PARSER_REGISTRY
                parse_fn = PARSER_REGISTRY["pdf"]
                doc = parse_fn(f)
                assert doc.source_path == str(f)
                assert "error" in doc.metadata
        finally:
            if saved is not None:
                sys.modules["fitz"] = saved

    def test_full_parse_requires_pymupdf(self, tmp_path):
        fitz = pytest.importorskip("fitz", reason="PyMuPDF not installed")
        # Create a minimal in-memory PDF using fitz
        pdf_path = tmp_path / "sample.pdf"
        doc = fitz.open()
        page = doc.new_page()
        page.insert_text((72, 72), "Hello from PDF")
        doc.save(str(pdf_path))
        doc.close()

        from homie_core.rag.parsers.pdf import parse_pdf
        result = parse_pdf(pdf_path)
        assert result.metadata["format"] == "pdf"
        assert result.metadata["pages"] == 1
        assert result.source_path == str(pdf_path)
        assert any("Hello from PDF" in b.content for b in result.text_blocks)

    def test_full_parse_page_numbers(self, tmp_path):
        fitz = pytest.importorskip("fitz", reason="PyMuPDF not installed")
        pdf_path = tmp_path / "two_pages.pdf"
        doc = fitz.open()
        for i in range(2):
            page = doc.new_page()
            page.insert_text((72, 72), f"Page {i + 1} content")
        doc.save(str(pdf_path))
        doc.close()

        from homie_core.rag.parsers.pdf import parse_pdf
        result = parse_pdf(pdf_path)
        assert result.metadata["pages"] == 2
        pages_found = {b.page for b in result.text_blocks}
        assert pages_found == {1, 2}


# ---------------------------------------------------------------------------
# DOCX Parser
# ---------------------------------------------------------------------------

class TestDocxParser:
    def test_registers_in_registry(self):
        import homie_core.rag.parsers.docx  # noqa: F401
        from homie_core.rag.parsers import PARSER_REGISTRY
        assert "docx" in PARSER_REGISTRY

    def test_fallback_missing_dep(self, tmp_path):
        f = tmp_path / "test.docx"
        f.write_bytes(b"PK fake docx")
        with patch.dict(sys.modules, {"docx": None}):
            import importlib
            import homie_core.rag.parsers.docx as docx_mod
            importlib.reload(docx_mod)
            from homie_core.rag.parsers import PARSER_REGISTRY
            parse_fn = PARSER_REGISTRY["docx"]
            doc = parse_fn(f)
            assert doc.source_path == str(f)
            assert "error" in doc.metadata

    def test_full_parse(self, tmp_path):
        pytest.importorskip("docx", reason="python-docx not installed")
        from docx import Document as DocxDocument

        docx_path = tmp_path / "sample.docx"
        d = DocxDocument()
        d.add_heading("Test Heading", level=1)
        d.add_paragraph("A sample paragraph.")
        d.save(str(docx_path))

        from homie_core.rag.parsers.docx import parse_docx
        result = parse_docx(docx_path)
        assert result.metadata["format"] == "docx"
        assert result.source_path == str(docx_path)
        headings = [b for b in result.text_blocks if b.block_type == "heading"]
        paragraphs = [b for b in result.text_blocks if b.block_type == "paragraph"]
        assert len(headings) >= 1
        assert headings[0].content == "Test Heading"
        assert headings[0].level == 1
        assert any("sample paragraph" in p.content for p in paragraphs)

    def test_full_parse_table(self, tmp_path):
        pytest.importorskip("docx", reason="python-docx not installed")
        from docx import Document as DocxDocument

        docx_path = tmp_path / "table.docx"
        d = DocxDocument()
        table = d.add_table(rows=2, cols=2)
        table.rows[0].cells[0].text = "Name"
        table.rows[0].cells[1].text = "Age"
        table.rows[1].cells[0].text = "Alice"
        table.rows[1].cells[1].text = "30"
        d.save(str(docx_path))

        from homie_core.rag.parsers.docx import parse_docx
        result = parse_docx(docx_path)
        assert len(result.tables) == 1
        assert result.tables[0].headers == ["Name", "Age"]
        assert result.tables[0].rows == [["Alice", "30"]]


# ---------------------------------------------------------------------------
# XLSX Parser
# ---------------------------------------------------------------------------

class TestXlsxParser:
    def test_registers_in_registry(self):
        import homie_core.rag.parsers.xlsx  # noqa: F401
        from homie_core.rag.parsers import PARSER_REGISTRY
        assert "xlsx" in PARSER_REGISTRY

    def test_fallback_missing_dep(self, tmp_path):
        f = tmp_path / "test.xlsx"
        f.write_bytes(b"PK fake xlsx")
        with patch.dict(sys.modules, {"openpyxl": None}):
            import importlib
            import homie_core.rag.parsers.xlsx as xlsx_mod
            importlib.reload(xlsx_mod)
            from homie_core.rag.parsers import PARSER_REGISTRY
            parse_fn = PARSER_REGISTRY["xlsx"]
            doc = parse_fn(f)
            assert doc.source_path == str(f)
            assert "error" in doc.metadata

    def test_full_parse(self, tmp_path):
        openpyxl = pytest.importorskip("openpyxl", reason="openpyxl not installed")

        xlsx_path = tmp_path / "sample.xlsx"
        wb = openpyxl.Workbook()
        ws = wb.active
        ws.title = "Sheet1"
        ws.append(["Name", "Score"])
        ws.append(["Alice", 95])
        ws.append(["Bob", 87])
        wb.save(str(xlsx_path))

        from homie_core.rag.parsers.xlsx import parse_xlsx
        result = parse_xlsx(xlsx_path)
        assert result.metadata["format"] == "xlsx"
        assert result.source_path == str(xlsx_path)
        assert len(result.tables) == 1
        assert result.tables[0].headers == ["Name", "Score"]
        assert result.tables[0].caption == "Sheet1"
        assert len(result.tables[0].rows) == 2

    def test_full_parse_sheet_count(self, tmp_path):
        openpyxl = pytest.importorskip("openpyxl", reason="openpyxl not installed")

        xlsx_path = tmp_path / "multi.xlsx"
        wb = openpyxl.Workbook()
        wb.active.title = "Alpha"
        wb.active.append(["x"])
        ws2 = wb.create_sheet("Beta")
        ws2.append(["y"])
        wb.save(str(xlsx_path))

        from homie_core.rag.parsers.xlsx import parse_xlsx
        result = parse_xlsx(xlsx_path)
        assert result.metadata["sheets"] == 2
        assert len(result.tables) == 2


# ---------------------------------------------------------------------------
# PPTX Parser
# ---------------------------------------------------------------------------

class TestPptxParser:
    def test_registers_in_registry(self):
        import homie_core.rag.parsers.pptx  # noqa: F401
        from homie_core.rag.parsers import PARSER_REGISTRY
        assert "pptx" in PARSER_REGISTRY

    def test_fallback_missing_dep(self, tmp_path):
        f = tmp_path / "test.pptx"
        f.write_bytes(b"PK fake pptx")
        with patch.dict(sys.modules, {"pptx": None}):
            import importlib
            import homie_core.rag.parsers.pptx as pptx_mod
            importlib.reload(pptx_mod)
            from homie_core.rag.parsers import PARSER_REGISTRY
            parse_fn = PARSER_REGISTRY["pptx"]
            doc = parse_fn(f)
            assert doc.source_path == str(f)
            assert "error" in doc.metadata

    def test_full_parse(self, tmp_path):
        pptx_lib = pytest.importorskip("pptx", reason="python-pptx not installed")
        from pptx import Presentation
        from pptx.util import Inches, Pt

        pptx_path = tmp_path / "sample.pptx"
        prs = Presentation()
        slide_layout = prs.slide_layouts[5]  # blank layout
        slide = prs.slides.add_slide(slide_layout)
        txBox = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(4), Inches(1))
        txBox.text_frame.text = "Slide one content"
        prs.save(str(pptx_path))

        from homie_core.rag.parsers.pptx import parse_pptx
        result = parse_pptx(pptx_path)
        assert result.metadata["format"] == "pptx"
        assert result.metadata["slides"] == 1
        assert result.source_path == str(pptx_path)
        assert any("Slide one content" in b.content for b in result.text_blocks)

    def test_full_parse_page_numbers(self, tmp_path):
        pytest.importorskip("pptx", reason="python-pptx not installed")
        from pptx import Presentation
        from pptx.util import Inches

        pptx_path = tmp_path / "two_slides.pptx"
        prs = Presentation()
        layout = prs.slide_layouts[5]
        for i in range(2):
            slide = prs.slides.add_slide(layout)
            tb = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(4), Inches(1))
            tb.text_frame.text = f"Slide {i + 1}"
        prs.save(str(pptx_path))

        from homie_core.rag.parsers.pptx import parse_pptx
        result = parse_pptx(pptx_path)
        pages = [b.page for b in result.text_blocks]
        assert 1 in pages
        assert 2 in pages


# ---------------------------------------------------------------------------
# HTML Parser
# ---------------------------------------------------------------------------

class TestHtmlParser:
    def test_registers_in_registry(self):
        import homie_core.rag.parsers.html  # noqa: F401
        from homie_core.rag.parsers import PARSER_REGISTRY
        assert "html" in PARSER_REGISTRY

    def test_basic_html_no_deps(self, tmp_path):
        """HTML parser works with stdlib only (regex fallback)."""
        f = tmp_path / "page.html"
        f.write_text("<html><body><p>Hello World</p></body></html>", encoding="utf-8")

        # Force fallback by blocking both trafilatura and bs4
        with patch.dict(sys.modules, {"trafilatura": None, "bs4": None}):
            import importlib
            import homie_core.rag.parsers.html as html_mod
            importlib.reload(html_mod)
            from homie_core.rag.parsers import PARSER_REGISTRY
            parse_fn = PARSER_REGISTRY["html"]
            doc = parse_fn(f)
            assert doc.source_path == str(f)
            assert doc.metadata["format"] == "html"
            assert "Hello World" in doc.full_text

    def test_bs4_fallback(self, tmp_path):
        pytest.importorskip("bs4", reason="beautifulsoup4 not installed")
        f = tmp_path / "page.html"
        f.write_text("<html><body><h1>Title</h1><p>Body text.</p></body></html>", encoding="utf-8")

        with patch.dict(sys.modules, {"trafilatura": None}):
            import importlib
            import homie_core.rag.parsers.html as html_mod
            importlib.reload(html_mod)
            from homie_core.rag.parsers import PARSER_REGISTRY
            parse_fn = PARSER_REGISTRY["html"]
            doc = parse_fn(f)
            assert "Title" in doc.full_text
            assert "Body text" in doc.full_text

    def test_metadata_format(self, tmp_path):
        f = tmp_path / "page.html"
        f.write_text("<p>test</p>", encoding="utf-8")
        from homie_core.rag.parsers.html import parse_html
        doc = parse_html(f)
        assert doc.metadata["format"] == "html"


# ---------------------------------------------------------------------------
# Image Parser
# ---------------------------------------------------------------------------

class TestImageParser:
    def test_registers_in_registry(self):
        import homie_core.rag.parsers.image  # noqa: F401
        from homie_core.rag.parsers import PARSER_REGISTRY
        assert "image" in PARSER_REGISTRY

    def test_fallback_missing_dep(self, tmp_path):
        f = tmp_path / "test.png"
        f.write_bytes(b"\x89PNG\r\n fake image bytes")
        with patch.dict(sys.modules, {"easyocr": None}):
            import importlib
            import homie_core.rag.parsers.image as image_mod
            importlib.reload(image_mod)
            from homie_core.rag.parsers import PARSER_REGISTRY
            parse_fn = PARSER_REGISTRY["image"]
            doc = parse_fn(f)
            assert doc.source_path == str(f)
            assert doc.metadata.get("format") == "image"
            assert "error" in doc.metadata

    def test_full_parse(self, tmp_path):
        pytest.importorskip("easyocr", reason="easyocr not installed")
        pytest.importorskip("PIL", reason="Pillow not installed")
        from PIL import Image, ImageDraw, ImageFont

        img_path = tmp_path / "text_image.png"
        img = Image.new("RGB", (200, 60), color=(255, 255, 255))
        draw = ImageDraw.Draw(img)
        draw.text((10, 10), "Hello OCR", fill=(0, 0, 0))
        img.save(str(img_path))

        from homie_core.rag.parsers.image import parse_image
        result = parse_image(img_path)
        assert result.metadata["format"] == "image"
        assert result.source_path == str(img_path)
        assert "ocr_regions" in result.metadata


# ---------------------------------------------------------------------------
# Email Parser
# ---------------------------------------------------------------------------

class TestEmailParser:
    def test_registers_in_registry(self):
        import homie_core.rag.parsers.email_parser  # noqa: F401
        from homie_core.rag.parsers import PARSER_REGISTRY
        assert "email" in PARSER_REGISTRY

    def _write_eml(self, path: Path, subject: str, body: str, sender: str = "alice@example.com",
                   to: str = "bob@example.com") -> None:
        content = (
            f"From: {sender}\r\n"
            f"To: {to}\r\n"
            f"Date: Sat, 21 Mar 2026 10:00:00 +0000\r\n"
            f"Subject: {subject}\r\n"
            f"Content-Type: text/plain; charset=utf-8\r\n"
            f"\r\n"
            f"{body}\r\n"
        )
        path.write_bytes(content.encode("utf-8"))

    def test_basic_plain_text_email(self, tmp_path):
        eml = tmp_path / "message.eml"
        self._write_eml(eml, subject="Hello", body="This is the body.")
        from homie_core.rag.parsers.email_parser import parse_email
        doc = parse_email(eml)
        assert doc.metadata["format"] == "email"
        assert doc.metadata["subject"] == "Hello"
        assert doc.metadata["from"] == "alice@example.com"
        assert doc.metadata["to"] == "bob@example.com"
        assert doc.source_path == str(eml)
        assert any("This is the body" in b.content for b in doc.text_blocks)

    def test_header_block_type(self, tmp_path):
        eml = tmp_path / "msg.eml"
        self._write_eml(eml, subject="Test subject", body="Body content")
        from homie_core.rag.parsers.email_parser import parse_email
        doc = parse_email(eml)
        header_blocks = [b for b in doc.text_blocks if b.block_type == "heading"]
        assert len(header_blocks) >= 1
        assert "Subject: Test subject" in header_blocks[0].content

    def test_attachments_list_empty(self, tmp_path):
        eml = tmp_path / "no_attach.eml"
        self._write_eml(eml, subject="No attach", body="No files here.")
        from homie_core.rag.parsers.email_parser import parse_email
        doc = parse_email(eml)
        assert doc.metadata["attachments"] == []

    def test_multipart_email(self, tmp_path):
        content = (
            "From: sender@example.com\r\n"
            "To: recv@example.com\r\n"
            "Subject: Multipart\r\n"
            "MIME-Version: 1.0\r\n"
            "Content-Type: multipart/mixed; boundary=\"boundary42\"\r\n"
            "\r\n"
            "--boundary42\r\n"
            "Content-Type: text/plain; charset=utf-8\r\n"
            "\r\n"
            "Plain text part\r\n"
            "--boundary42--\r\n"
        )
        eml = tmp_path / "multi.eml"
        eml.write_bytes(content.encode("utf-8"))
        from homie_core.rag.parsers.email_parser import parse_email
        doc = parse_email(eml)
        assert doc.metadata["subject"] == "Multipart"
        assert any("Plain text part" in b.content for b in doc.text_blocks)


# ---------------------------------------------------------------------------
# Code Parser
# ---------------------------------------------------------------------------

class TestCodeParser:
    def test_registers_in_registry(self):
        import homie_core.rag.parsers.code  # noqa: F401
        from homie_core.rag.parsers import PARSER_REGISTRY
        assert "code" in PARSER_REGISTRY

    def test_python_file(self, tmp_path):
        f = tmp_path / "script.py"
        f.write_text("def hello():\n    return 'world'\n", encoding="utf-8")
        from homie_core.rag.parsers.code import parse_code
        doc = parse_code(f)
        assert doc.metadata["format"] == "code"
        assert doc.metadata["language"] == "python"
        assert doc.source_path == str(f)
        assert len(doc.text_blocks) == 1
        assert doc.text_blocks[0].block_type == "code"
        assert doc.text_blocks[0].language == "python"

    def test_javascript_file(self, tmp_path):
        f = tmp_path / "app.js"
        f.write_text("console.log('hi');", encoding="utf-8")
        from homie_core.rag.parsers.code import parse_code
        doc = parse_code(f)
        assert doc.metadata["language"] == "javascript"

    def test_unknown_extension(self, tmp_path):
        f = tmp_path / "data.xyz"
        f.write_text("some content", encoding="utf-8")
        from homie_core.rag.parsers.code import parse_code
        doc = parse_code(f)
        assert doc.metadata["language"] == "unknown"

    def test_line_count(self, tmp_path):
        f = tmp_path / "three.py"
        f.write_text("a = 1\nb = 2\nc = 3\n", encoding="utf-8")
        from homie_core.rag.parsers.code import parse_code
        doc = parse_code(f)
        assert doc.metadata["lines"] == 4  # 3 newlines + 1

    @pytest.mark.parametrize("ext,lang", [
        (".ts", "typescript"),
        (".go", "go"),
        (".rs", "rust"),
        (".java", "java"),
        (".rb", "ruby"),
        (".sh", "bash"),
        (".sql", "sql"),
        (".php", "php"),
    ])
    def test_language_mapping(self, tmp_path, ext, lang):
        f = tmp_path / f"file{ext}"
        f.write_text("code here", encoding="utf-8")
        from homie_core.rag.parsers.code import parse_code
        doc = parse_code(f)
        assert doc.metadata["language"] == lang

    def test_encoding_error_handled(self, tmp_path):
        f = tmp_path / "latin.py"
        f.write_bytes(b"# caf\xe9\nx = 1\n")
        from homie_core.rag.parsers.code import parse_code
        doc = parse_code(f)  # should not raise
        assert doc.metadata["language"] == "python"


# ---------------------------------------------------------------------------
# Epub Parser
# ---------------------------------------------------------------------------

class TestEpubParser:
    def test_registers_in_registry(self):
        import homie_core.rag.parsers.epub  # noqa: F401
        from homie_core.rag.parsers import PARSER_REGISTRY
        assert "epub" in PARSER_REGISTRY

    def test_fallback_missing_dep(self, tmp_path):
        f = tmp_path / "test.epub"
        f.write_bytes(b"PK fake epub")
        with patch.dict(sys.modules, {"ebooklib": None}):
            import importlib
            import homie_core.rag.parsers.epub as epub_mod
            importlib.reload(epub_mod)
            from homie_core.rag.parsers import PARSER_REGISTRY
            parse_fn = PARSER_REGISTRY["epub"]
            doc = parse_fn(f)
            assert doc.source_path == str(f)
            assert "error" in doc.metadata

    def test_full_parse(self, tmp_path):
        pytest.importorskip("ebooklib", reason="ebooklib not installed")
        from ebooklib import epub

        epub_path = tmp_path / "sample.epub"
        book = epub.EpubBook()
        book.set_title("Test Book")
        book.set_language("en")

        chapter = epub.EpubHtml(title="Chapter 1", file_name="chap1.xhtml", lang="en")
        chapter.content = b"<html><body><p>Chapter one content</p></body></html>"
        book.add_item(chapter)
        book.spine = ["nav", chapter]
        book.add_item(epub.EpubNcx())
        book.add_item(epub.EpubNav())
        epub.write_epub(str(epub_path), book)

        from homie_core.rag.parsers.epub import parse_epub
        result = parse_epub(epub_path)
        assert result.metadata["format"] == "epub"
        assert result.metadata["title"] == "Test Book"
        assert result.source_path == str(epub_path)
        assert any("Chapter one content" in b.content for b in result.text_blocks)


# ---------------------------------------------------------------------------
# Registry completeness: all 9 parsers registered after importing modules
# ---------------------------------------------------------------------------

class TestAllParsersRegistered:
    def test_all_binary_parsers_in_registry(self):
        # Import all parser modules to trigger registration
        import homie_core.rag.parsers.pdf  # noqa: F401
        import homie_core.rag.parsers.docx  # noqa: F401
        import homie_core.rag.parsers.xlsx  # noqa: F401
        import homie_core.rag.parsers.pptx  # noqa: F401
        import homie_core.rag.parsers.html  # noqa: F401
        import homie_core.rag.parsers.image  # noqa: F401
        import homie_core.rag.parsers.email_parser  # noqa: F401
        import homie_core.rag.parsers.code  # noqa: F401
        import homie_core.rag.parsers.epub  # noqa: F401

        from homie_core.rag.parsers import PARSER_REGISTRY
        expected = {"pdf", "docx", "xlsx", "pptx", "html", "image", "email", "code", "epub"}
        assert expected.issubset(set(PARSER_REGISTRY.keys()))

    def test_all_parsers_return_parsed_document(self, tmp_path):
        """Each parser returns a ParsedDocument even for unrecognised/dummy files."""
        import homie_core.rag.parsers.pdf  # noqa: F401
        import homie_core.rag.parsers.docx  # noqa: F401
        import homie_core.rag.parsers.xlsx  # noqa: F401
        import homie_core.rag.parsers.pptx  # noqa: F401
        import homie_core.rag.parsers.html  # noqa: F401
        import homie_core.rag.parsers.image  # noqa: F401
        import homie_core.rag.parsers.email_parser  # noqa: F401
        import homie_core.rag.parsers.code  # noqa: F401
        import homie_core.rag.parsers.epub  # noqa: F401

        from homie_core.rag.parsers import PARSER_REGISTRY, ParsedDocument

        dummy_files = {
            "pdf": tmp_path / "d.pdf",
            "docx": tmp_path / "d.docx",
            "xlsx": tmp_path / "d.xlsx",
            "pptx": tmp_path / "d.pptx",
            "html": tmp_path / "d.html",
            "image": tmp_path / "d.png",
            "email": tmp_path / "d.eml",
            "code": tmp_path / "d.py",
            "epub": tmp_path / "d.epub",
        }
        # Write minimal content so parsers that do plain reads don't blow up
        for fmt, fp in dummy_files.items():
            if fmt == "html":
                fp.write_text("<p>test</p>", encoding="utf-8")
            elif fmt == "email":
                fp.write_bytes(
                    b"From: a@b.com\r\nSubject: t\r\n\r\nbody\r\n"
                )
            elif fmt == "code":
                fp.write_text("x = 1", encoding="utf-8")
            else:
                fp.write_bytes(b"fake binary content")

        for fmt, fp in dummy_files.items():
            parse_fn = PARSER_REGISTRY[fmt]
            try:
                result = parse_fn(fp)
                assert isinstance(result, ParsedDocument), (
                    f"Parser '{fmt}' did not return ParsedDocument, got {type(result)}"
                )
                assert result.source_path == str(fp), (
                    f"Parser '{fmt}' did not set source_path correctly"
                )
            except Exception as exc:
                # Parsers with missing deps return gracefully — any exception is a bug
                pytest.fail(f"Parser '{fmt}' raised an unexpected exception: {exc}")
